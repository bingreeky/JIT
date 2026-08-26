"""
Slimmed-down MarkdownConverter for converting documents to markdown text.

Adapted from Flash-Searcher's mm_tools_utils.py. Keeps only the converters
needed for our benchmarks: PlainText, HTML, PDF, DOCX, XLSX, PPTX, ZIP.

Dropped: Wikipedia, YouTube, WAV/MP3/Image/Media converters and their
dependencies (puremagic, pydub, speech_recognition, youtube_transcript_api).
"""

import json
import logging
import os
import re
import shutil
import tempfile
import traceback
import zipfile
from typing import Any, Dict, List, Optional, Union

import markdownify
from bs4 import BeautifulSoup


# pdfminer logs a WARNING *per font per page* for any FontDescriptor without a
# usable FontBBox ("Could not get FontBBox from font descriptor because None
# cannot be parsed as 4 floats"). It is cosmetic -- pdfminer substitutes a
# default bbox and text extraction is unaffected -- but subsetted/Type3 fonts
# make it fire hundreds of times for a single attachment, which buries the
# agent trajectory in the run logs. Errors still surface.
logging.getLogger("pdfminer").setLevel(logging.ERROR)


# ── Custom Markdownify ────────────────────────────────────────────────

class _CustomMarkdownify(markdownify.MarkdownConverter):
    """Enhanced HTML-to-Markdown converter."""

    def convert_hn(self, n, el, text, convert_as_inline):
        if not text.strip():
            return ""
        style = self.options["heading_style"]
        if style == markdownify.UNDERLINED and n <= 2:
            line = "=" if n == 1 else "-"
            return self.indent("\n" + text + "\n" + line * len(text) + "\n")
        hashes = "#" * n
        return self.indent("\n" + hashes + " " + text.strip() + "\n\n")

    def convert_a(self, el, text, convert_as_inline):
        prefix, suffix, text = markdownify.chomp(text)
        if not text:
            return ""
        href = el.get("href", "")
        title = el.get("title", "")
        if title:
            return f"{prefix}[{text}]({href} \"{title}\"){suffix}"
        return f"{prefix}[{text}]({href}){suffix}" if href else text

    def convert_img(self, el, text, convert_as_inline):
        alt = el.get("alt", "") or ""
        src = el.get("src", "") or ""
        title = el.get("title", "") or ""
        if not src:
            return alt
        # Skip data URIs (too long)
        if src.startswith("data:"):
            return f"![{alt}](data:...)"
        if title:
            return f"![{alt}]({src} \"{title}\")"
        return f"![{alt}]({src})"


def _md(html: str, **options) -> str:
    return _CustomMarkdownify(**options).convert(html)


# ── Converter Base ────────────────────────────────────────────────────

class DocumentConverterResult:
    """Result of document conversion."""

    def __init__(self, title: Optional[str] = None, text_content: str = ""):
        self.title = title
        self.text_content = text_content


class DocumentConverter:
    """Abstract base for document converters."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        raise NotImplementedError


# ── Converters ────────────────────────────────────────────────────────

class PlainTextConverter(DocumentConverter):
    """Read plain text files."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        # Handle common text/code extensions
        text_exts = {
            ".txt", ".md", ".csv", ".tsv", ".json", ".jsonl", ".xml",
            ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".log",
            ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
            ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".zsh",
            ".css", ".scss", ".less", ".sql", ".r", ".m", ".tex",
            ".srt", ".ass", ".vtt",
        }
        if ext not in text_exts:
            return None
        try:
            with open(local_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=content,
            )
        except Exception:
            return None


class HtmlConverter(DocumentConverter):
    """Convert HTML files to Markdown."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        if ext not in (".html", ".htm"):
            return None
        try:
            with open(local_path, "r", encoding="utf-8", errors="replace") as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, "html.parser")

            # Remove scripts and styles
            for tag in soup.find_all(["script", "style"]):
                tag.decompose()

            title = ""
            title_tag = soup.find("title")
            if title_tag:
                title = title_tag.get_text(strip=True)

            body = soup.find("body") or soup
            md_content = _md(str(body), heading_style="ATX", strip=["img"])

            return DocumentConverterResult(title=title, text_content=md_content.strip())
        except Exception:
            return None


class PdfConverter(DocumentConverter):
    """Convert PDF files to text."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        if ext != ".pdf":
            return None
        try:
            import pdfminer.high_level

            text = pdfminer.high_level.extract_text(local_path)
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=text.strip(),
            )
        except Exception as e:
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error reading PDF: {e}",
            )


class DocxConverter(DocumentConverter):
    """Convert DOCX files to Markdown."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        if ext != ".docx":
            return None
        try:
            import mammoth

            with open(local_path, "rb") as f:
                result = mammoth.convert_to_html(f)
            md_content = _md(result.value, heading_style="ATX")
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=md_content.strip(),
            )
        except Exception as e:
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error reading DOCX: {e}",
            )


class XlsxConverter(DocumentConverter):
    """Convert Excel files to Markdown tables."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        if ext not in (".xlsx", ".xls"):
            return None
        try:
            import pandas as pd

            sheets = pd.read_excel(local_path, sheet_name=None, engine="openpyxl")
            parts = []
            for sheet_name, df in sheets.items():
                parts.append(f"## Sheet: {sheet_name}")
                # Limit rows for large spreadsheets
                if len(df) > 100:
                    parts.append(f"(Showing first 100 of {len(df)} rows)")
                    df = df.head(100)
                html = df.to_html(index=False, na_rep="")
                md = _md(html, heading_style="ATX")
                parts.append(md)
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content="\n\n".join(parts).strip(),
            )
        except Exception as e:
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error reading Excel: {e}",
            )


class PptxConverter(DocumentConverter):
    """Convert PowerPoint files to text."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        if ext not in (".pptx", ".ppt"):
            return None
        try:
            from pptx import Presentation

            prs = Presentation(local_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = self._process_slide(slide)
                if slide_text.strip():
                    parts.append(f"## Slide {i}\n{slide_text}")
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content="\n\n".join(parts).strip(),
            )
        except Exception as e:
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error reading PPTX: {e}",
            )

    @staticmethod
    def _process_slide(slide) -> str:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    texts.append(f"| {row_text} |")
        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Speaker Notes]: {notes}")
        return "\n".join(texts)


class ZipConverter(DocumentConverter):
    """List contents of ZIP files."""

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        ext = os.path.splitext(local_path)[1].lower()
        if ext != ".zip":
            return None
        try:
            with zipfile.ZipFile(local_path, "r") as zf:
                file_list = zf.namelist()
            content = f"ZIP archive contents ({len(file_list)} files):\n"
            for f in file_list[:50]:
                content += f"  - {f}\n"
            if len(file_list) > 50:
                content += f"  ... and {len(file_list) - 50} more files\n"
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=content.strip(),
            )
        except Exception as e:
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error reading ZIP: {e}",
            )


# ── Main Converter ────────────────────────────────────────────────────

class MarkdownConverter:
    """Converts documents to Markdown text.

    Dispatches to registered converters based on file extension.
    """

    def __init__(self):
        self._converters: List[DocumentConverter] = []
        self._register_defaults()

    def _register_defaults(self) -> None:
        """Register default converters (order matters — last registered = tried first)."""
        self._converters = [
            PlainTextConverter(),
            HtmlConverter(),
            DocxConverter(),
            XlsxConverter(),
            PptxConverter(),
            ZipConverter(),
            PdfConverter(),
        ]

    def convert(self, local_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        """Convert a local file to Markdown text.

        Tries each converter in reverse order (last registered = highest priority).
        """
        if not os.path.isfile(local_path):
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error: file not found: {local_path}",
            )

        # Try converters in reverse order (last = highest priority)
        for converter in reversed(self._converters):
            try:
                result = converter.convert(local_path, **kwargs)
                if result is not None:
                    return self._normalize(result)
            except Exception:
                continue

        # Fallback: try reading as plain text
        try:
            with open(local_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(100000)
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=content,
            )
        except Exception:
            return DocumentConverterResult(
                title=os.path.basename(local_path),
                text_content=f"Error: could not read file: {local_path}",
            )

    @staticmethod
    def _normalize(result: DocumentConverterResult) -> DocumentConverterResult:
        """Clean up whitespace."""
        if result.text_content:
            # Remove trailing whitespace on each line
            lines = [line.rstrip() for line in result.text_content.split("\n")]
            # Collapse 3+ consecutive blank lines to 2
            text = "\n".join(lines)
            text = re.sub(r"\n{3,}", "\n\n", text)
            result.text_content = text.strip()
        return result
